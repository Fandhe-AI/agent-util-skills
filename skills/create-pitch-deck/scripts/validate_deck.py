#!/usr/bin/env python3
"""validate_deck.py — 生成 PPTX の機械検証。

役割と境界:
- build_deck.py（または手直しした pptx）の出力が SKILL.md の契約
  （はみ出しなし / 日本語フォント fallback / 前提と解釈・フィードバック観点の
  必須配置）を満たすかを、生成された .pptx ファイルそのものから検証する。
  create-html-report/scripts/validate_report.py と同じく、spec ではなく
  「出力された成果物」を検証対象にする（spec を信頼しない）。
- 検証のみを行い、ファイルの修正は行わない。

使い方:
    python3 validate_deck.py <output.pptx>

終了コード: 全チェック PASS で 0、1 件でも FAIL で 1。
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
except ImportError:
    print(
        "python-pptx が見つからない。venv へ `pip install python-pptx` を"
        "実行してから再実行すること。",
        file=sys.stderr,
    )
    sys.exit(1)

MIN_SLIDES = 8
MAX_SLIDES = 14
FEEDBACK_ITEM_MIN = 3
FEEDBACK_ITEM_MAX = 5
NUMBERED_ITEM_RE = re.compile(r"^\s*\d+\.\s*\S")


def iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from shape.shapes


def check_slide_count(prs, failures):
    n = len(prs.slides)
    if not (MIN_SLIDES <= n <= MAX_SLIDES):
        failures.append(
            f"スライド枚数が範囲外: {n}枚（想定 {MIN_SLIDES}〜{MAX_SLIDES}枚 = 10枚前後）"
        )


def check_bounds(prs, failures):
    sw, sh = prs.slide_width, prs.slide_height
    for i, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide):
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
            if left is None or top is None or width is None or height is None:
                continue
            if left < 0 or top < 0:
                failures.append(f"slide {i}: shape '{shape.shape_id}' が負の座標 (left={left}, top={top})")
            if left + width > sw:
                failures.append(
                    f"slide {i}: shape '{shape.shape_id}' がスライド右端をはみ出す "
                    f"(left+width={left + width} > slide_width={sw})"
                )
            if top + height > sh:
                failures.append(
                    f"slide {i}: shape '{shape.shape_id}' がスライド下端をはみ出す "
                    f"(top+height={top + height} > slide_height={sh})"
                )


def check_fonts(prs, failures):
    for i, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide):
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rpr = run._r.find(qn("a:rPr"))  # noqa: SLF001
                    if rpr is None:
                        failures.append(f"slide {i}: 文字列 '{run.text[:20]}' に rPr が無くフォント未指定")
                        continue
                    latin = rpr.find(qn("a:latin"))
                    ea = rpr.find(qn("a:ea"))
                    if latin is None or not latin.get("typeface"):
                        failures.append(f"slide {i}: 文字列 '{run.text[:20]}' に a:latin typeface が未設定")
                    if ea is None or not ea.get("typeface"):
                        failures.append(
                            f"slide {i}: 文字列 '{run.text[:20]}' に a:ea typeface が未設定"
                            "（Google スライド取込時の日本語フォールバックに必須）"
                        )


def slide_text(slide) -> str:
    parts = []
    for shape in iter_shapes(slide):
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return "\n".join(parts)


def check_premise_slide(prs, failures):
    slides = list(prs.slides)
    if len(slides) < 2:
        failures.append("スライドが2枚未満で「前提と解釈」スライドの位置を検証できない")
        return
    text = slide_text(slides[1])
    if "前提" not in text:
        failures.append("2枚目（前提と解釈スライド）に「前提」という語が含まれない")


def check_feedback_slide(prs, failures):
    slides = list(prs.slides)
    if not slides:
        failures.append("スライドが0枚で「フィードバック観点」スライドを検証できない")
        return
    last = slides[-1]
    text = slide_text(last)
    if "フィードバック" not in text:
        failures.append("最終スライドに「フィードバック」という語が含まれない")
        return
    items = [line for line in text.splitlines() if NUMBERED_ITEM_RE.match(line)]
    if not (FEEDBACK_ITEM_MIN <= len(items) <= FEEDBACK_ITEM_MAX):
        failures.append(
            "最終スライドの番号付きフィードバック観点が"
            f"{FEEDBACK_ITEM_MIN}〜{FEEDBACK_ITEM_MAX}件の範囲外（検出 {len(items)}件）"
        )


def main() -> int:
    if len(sys.argv) != 2:
        print("使い方: python3 validate_deck.py <output.pptx>", file=sys.stderr)
        return 1

    pptx_path = Path(sys.argv[1])
    if not pptx_path.is_file():
        print(f"FAIL: ファイルが存在しない: {pptx_path}")
        return 1

    prs = Presentation(str(pptx_path))
    failures: list[str] = []

    check_slide_count(prs, failures)
    check_bounds(prs, failures)
    check_fonts(prs, failures)
    check_premise_slide(prs, failures)
    check_feedback_slide(prs, failures)

    if failures:
        print(f"FAIL: {len(failures)}件の問題を検出")
        for f in failures:
            print(f" - {f}")
        return 1

    print(f"PASS: {pptx_path} は全チェックを通過（{len(prs.slides)}枚）")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
