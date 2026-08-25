#!/usr/bin/env python3
"""deck spec (JSON) から 16:9 の企画提案スライド (PPTX) を生成する renderer。

役割: report spec -> HTML の create-html-report/scripts/render_report.py と対になる
存在。座標計算・フォント設定 (a:latin と a:ea の両方)・画像はめ込みは本スクリプトの
責務とし、呼び出し側 (SKILL.md の手順) は spec の内容判断にのみ集中できるようにする。

スライド構成（PO 承認会向け）: 前半 = 課題・解決アプローチ・スコープ・勝ち筋、
後半 = 画面と操作の流れ（screen_flow, 2〜4枚。create-design-doc の storyboard/
screens 画像を取り込む）→ 検証計画 → 承認いただきたい事項＋確認事項。

依存: python-pptx (標準ライブラリ外。venv へインストールしてから実行する)
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Emu, Inches, Pt
except ImportError:
    print(
        "python-pptx が見つからない。venv を作成し `pip install python-pptx` を"
        "実行してから再実行すること（SKILL.md の Step 5 参照）。",
        file=sys.stderr,
    )
    sys.exit(1)

# 前半（固定・この順序）: 課題認識から勝ち筋まで
FRONT_HALF_ROLES = ["cover", "premise", "problem", "solution", "scope", "winning"]
# 後半固定末尾（screen_flow の直後、この順序・この2枚で終わること）
BACK_HALF_TAIL_ROLES = ["validation", "approval"]
SCREEN_FLOW_ROLE = "screen_flow"
SCREEN_FLOW_MIN = 2
SCREEN_FLOW_MAX = 4

APPROVAL_ITEM_MIN = 3
APPROVAL_ITEM_MAX = 5
APPROVAL_KINDS = {"承認", "確認"}
WINNING_LABELS = {"事実", "仮説"}

DEFAULT_BRAND = {
    "primary": "#1F3A93",
    "secondary": "#2C3E50",
    "accent": "#E08E45",
    "background": "#FFFFFF",
    "surface": "#F4F6F8",
    "text": "#1C1C1C",
    "muted": "#5B6470",
    "font_latin": "Arial",
    "font_ea": "Noto Sans JP",
}

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


class SpecError(Exception):
    """spec の構造・内容が契約を満たさない場合に送出する。"""


def load_spec(path: Path) -> dict:
    try:
        raw = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise SpecError(f"{path} が妥当な JSON でない: {exc}") from exc
    if not isinstance(raw, dict):
        raise SpecError("spec のルートは JSON object であること")
    return raw


def validate_spec(spec: dict) -> None:
    """generate 前の構造チェック。

    構成契約: 前半 [cover, premise, problem, solution, scope, winning]
    （この順・各1枚）→ screen_flow 2〜4枚（連続） → [validation, approval]
    （この順・各1枚、これで終わること）。
    """
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        raise SpecError("spec.slides は 1 件以上の配列であること")

    roles = [s.get("role") for s in slides if isinstance(s, dict)]

    prefix = roles[: len(FRONT_HALF_ROLES)]
    if prefix != FRONT_HALF_ROLES:
        raise SpecError(
            f"前半 {len(FRONT_HALF_ROLES)}枚は role={FRONT_HALF_ROLES} の順であること"
            f"（現在: {prefix}）"
        )

    idx = len(FRONT_HALF_ROLES)
    n_flow = 0
    while idx < len(roles) and roles[idx] == SCREEN_FLOW_ROLE:
        n_flow += 1
        idx += 1
    if not (SCREEN_FLOW_MIN <= n_flow <= SCREEN_FLOW_MAX):
        raise SpecError(
            f"role={SCREEN_FLOW_ROLE}（画面と操作の流れ）は"
            f"{SCREEN_FLOW_MIN}〜{SCREEN_FLOW_MAX}枚連続で配置すること"
            f"（現在 {n_flow}枚）。PO 承認会でシナリオごとに画面の使われ方を説明する枠"
        )

    suffix = roles[idx:]
    if suffix != BACK_HALF_TAIL_ROLES:
        raise SpecError(
            f"role={SCREEN_FLOW_ROLE} の直後は role={BACK_HALF_TAIL_ROLES} の順で"
            f"終わること（現在: {suffix}）"
        )

    for slide in slides:
        role = slide.get("role")
        title = slide.get("title")
        if not title or not isinstance(title, str):
            raise SpecError(f"role={role} に title (string) が必要")

        if role == "premise":
            _require_str_list(slide, "bullets", role, min_len=1)
        elif role in ("problem", "solution", "validation"):
            _require_str_list(slide, "bullets", role, min_len=1)
        elif role == "scope":
            _require_str_list(slide, "in_scope", role, min_len=1)
            _require_str_list(slide, "out_scope", role, min_len=1)
        elif role == "winning":
            items = slide.get("items")
            if not isinstance(items, list) or not items:
                raise SpecError("role=winning に items (1件以上) が必要")
            for item in items:
                if not isinstance(item, dict) or not item.get("text"):
                    raise SpecError("winning.items の各要素は text を持つこと")
                label = item.get("label")
                if label not in WINNING_LABELS:
                    raise SpecError(
                        "winning.items の各要素は label に "
                        f"{sorted(WINNING_LABELS)} のいずれかを指定すること。"
                        "'事実' は入力文書に記録された実測・調査結果に限る"
                        "（留保がある場合は text に併記する）。それ以外は '仮説'"
                    )
        elif role == SCREEN_FLOW_ROLE:
            narrative = slide.get("narrative")
            if not narrative or not isinstance(narrative, str):
                raise SpecError(
                    f"role={SCREEN_FLOW_ROLE} に narrative (string) が必要。"
                    "「この場面で・この画面が・こう使われる」を説明する文"
                )
            image = slide.get("image")
            if image is not None and not isinstance(image, str):
                raise SpecError(f"role={SCREEN_FLOW_ROLE}.image は string か null であること")
            if not image:
                note = slide.get("note")
                if not note or not isinstance(note, str):
                    raise SpecError(
                        f"role={SCREEN_FLOW_ROLE} で image が無い場合は note (string) が必要。"
                        "例: 'create-design-doc 未実行のためテキスト概略のみ'"
                    )
        elif role == "approval":
            items = slide.get("items")
            if not isinstance(items, list):
                raise SpecError("role=approval に items (配列) が必要")
            if not (APPROVAL_ITEM_MIN <= len(items) <= APPROVAL_ITEM_MAX):
                raise SpecError(
                    "role=approval の items は"
                    f"{APPROVAL_ITEM_MIN}〜{APPROVAL_ITEM_MAX}件であること"
                    f"（現在 {len(items)}件）"
                )
            for item in items:
                if not isinstance(item, dict) or not item.get("text"):
                    raise SpecError("approval.items の各要素は text を持つこと")
                if item.get("kind") not in APPROVAL_KINDS:
                    raise SpecError(
                        f"approval.items の各要素は kind に {sorted(APPROVAL_KINDS)} の"
                        "いずれかを指定すること"
                    )


def _require_str_list(slide: dict, key: str, role: str, min_len: int = 1) -> None:
    value = slide.get(key)
    if not isinstance(value, list) or len(value) < min_len:
        raise SpecError(f"role={role} に {key} ({min_len}件以上の配列) が必要")
    for v in value:
        if not isinstance(v, str) or not v.strip():
            raise SpecError(f"role={role}.{key} の各要素は非空文字列であること")


def hex_to_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    if len(value) != 6:
        raise SpecError(f"不正なカラーコード: {value}")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_run_font(run, brand: dict, size: int, bold: bool = False, color: str | None = None) -> None:
    """a:latin と a:ea の両方に typeface を設定する。

    python-pptx の Font.name は a:latin にしか書き込まないため、a:ea は
    lxml 経由で明示的に追加する（Google スライド取込時の日本語フォールバック対応）。
    """
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = brand["font_latin"]
    rpr = run.font._rPr  # noqa: SLF001 - python-pptx が公開 API を提供していないため
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        rpr.append(ea)
    ea.set("typeface", brand["font_ea"])
    run.font.color.rgb = hex_to_rgb(color or brand["text"])


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return box, tf


def add_paragraph_text(tf, text, brand, size, bold=False, color=None, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, brand, size, bold=bold, color=color)
    return p


def add_rect(slide, left, top, width, height, fill_hex):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))  # 1 = RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_picture_fit(slide, path, left, top, max_w, max_h):
    """画像を (max_w, max_h) の枠内にアスペクト比を保ったまま収めて配置する。

    Pillow に依存せず python-pptx 自身のネイティブサイズ計算だけで完結させる:
    まず width 基準で仮配置して自然な height を得る。height が枠を超える場合
    （storyboard.png のような縦長画像）は一度削除し height 基準で再配置する
    （結果の width は枠内に収まることが幅と高さの比の性質上保証される）。
    """
    pic = slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(max_w))
    if pic.height > Inches(max_h):
        pic._element.getparent().remove(pic._element)  # noqa: SLF001
        pic = slide.shapes.add_picture(path, Inches(left), Inches(top), height=Inches(max_h))
    # 枠内で中央寄せする
    pic.left = Inches(left) + (Inches(max_w) - pic.width) // 2
    pic.top = Inches(top) + (Inches(max_h) - pic.height) // 2
    return pic


def add_title_bar(slide, title, brand, subtitle=None):
    """cover 以外の全スライド共通のヘッダー（アクセントバー + タイトル）。"""
    add_rect(slide, 0, 0, SLIDE_W_IN, 0.12, brand["primary"])
    _, tf = add_textbox(slide, 0.6, 0.35, SLIDE_W_IN - 1.2, 0.9)
    add_paragraph_text(tf, title, brand, size=28, bold=True, color=brand["secondary"], first=True)
    if subtitle:
        _, sub_tf = add_textbox(slide, 0.6, 1.05, SLIDE_W_IN - 1.2, 0.4)
        add_paragraph_text(sub_tf, subtitle, brand, size=14, color=brand["muted"], first=True)


def add_footer(slide, index, total, brand, deck_title):
    _, tf = add_textbox(slide, 0.6, SLIDE_H_IN - 0.45, SLIDE_W_IN - 1.6, 0.35)
    add_paragraph_text(tf, deck_title, brand, size=9, color=brand["muted"], first=True)
    _, page_tf = add_textbox(slide, SLIDE_W_IN - 1.0, SLIDE_H_IN - 0.45, 0.6, 0.35)
    page_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    add_paragraph_text(page_tf, f"{index}/{total}", brand, size=9, color=brand["muted"], first=True)


def add_bullets(slide, bullets, brand, top, left=0.6, width=SLIDE_W_IN - 1.2, size=18, marker="—  "):
    _, tf = add_textbox(slide, left, top, width, SLIDE_H_IN - top - 0.7)
    for i, text in enumerate(bullets):
        add_paragraph_text(tf, f"{marker}{text}", brand, size=size, color=brand["text"], first=(i == 0))
        if i < len(bullets) - 1:
            tf.paragraphs[-1].space_after = Pt(10)


def build_cover(prs, slide_spec, brand, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, brand["primary"])
    add_rect(slide, 0, SLIDE_H_IN - 1.6, SLIDE_W_IN, 1.6, brand["secondary"])
    _, tf = add_textbox(slide, 0.9, 2.3, SLIDE_W_IN - 1.8, 1.8, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraph_text(tf, slide_spec["title"], brand, size=40, bold=True, color="#FFFFFF", first=True)
    subtitle = slide_spec.get("subtitle")
    if subtitle:
        _, sub_tf = add_textbox(slide, 0.9, 3.9, SLIDE_W_IN - 1.8, 0.7)
        add_paragraph_text(sub_tf, subtitle, brand, size=18, color="#F4F6F8", first=True)
    meta_bits = [b for b in (slide_spec.get("date"), slide_spec.get("meta")) if b]
    if meta_bits:
        _, meta_tf = add_textbox(slide, 0.9, SLIDE_H_IN - 1.15, SLIDE_W_IN - 1.8, 0.5)
        add_paragraph_text(meta_tf, "  |  ".join(meta_bits), brand, size=12, color="#D8DEE6", first=True)
    return slide


def build_premise(prs, slide_spec, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, slide_spec["title"], brand)
    add_bullets(slide, slide_spec["bullets"], brand, top=1.6)
    source_note = slide_spec.get("source_note")
    if source_note:
        _, tf = add_textbox(slide, 0.6, SLIDE_H_IN - 1.05, SLIDE_W_IN - 1.2, 0.4)
        add_paragraph_text(tf, f"出典: {source_note}", brand, size=11, color=brand["muted"], first=True)
    return slide


def build_bullets_slide(prs, slide_spec, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, slide_spec["title"], brand)
    add_bullets(slide, slide_spec["bullets"], brand, top=1.6)
    return slide


def build_scope(prs, slide_spec, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, slide_spec["title"], brand)
    col_w = (SLIDE_W_IN - 1.2 - 0.4) / 2
    add_rect(slide, 0.6, 1.6, col_w, SLIDE_H_IN - 2.3, brand["surface"])
    add_rect(slide, 0.6 + col_w + 0.4, 1.6, col_w, SLIDE_H_IN - 2.3, brand["surface"])
    _, in_head = add_textbox(slide, 0.8, 1.75, col_w - 0.4, 0.4)
    add_paragraph_text(in_head, "In Scope", brand, size=16, bold=True, color=brand["primary"], first=True)
    add_bullets(slide, slide_spec["in_scope"], brand, top=2.3, left=0.8, width=col_w - 0.4, size=15)
    _, out_head = add_textbox(slide, 0.8 + col_w + 0.4, 1.75, col_w - 0.4, 0.4)
    add_paragraph_text(out_head, "Out of Scope", brand, size=16, bold=True, color=brand["muted"], first=True)
    add_bullets(slide, slide_spec["out_scope"], brand, top=2.3, left=0.8 + col_w + 0.4, width=col_w - 0.4, size=15)
    return slide


def build_winning(prs, slide_spec, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, slide_spec["title"], brand)
    _, tf = add_textbox(slide, 0.6, 1.6, SLIDE_W_IN - 1.2, SLIDE_H_IN - 2.3)
    for i, item in enumerate(slide_spec["items"]):
        label = item["label"]
        badge_color = brand["primary"] if label == "事実" else brand["accent"]
        text = f"[{label}] {item['text']}"
        add_paragraph_text(tf, f"—  {text}", brand, size=18, color=badge_color if label == "仮説" else brand["text"], first=(i == 0))
        if i < len(slide_spec["items"]) - 1:
            tf.paragraphs[-1].space_after = Pt(10)
    return slide


def build_screen_flow(prs, slide_spec, brand):
    """「画面と操作の流れ」スライド。左に画面画像（storyboard/screens 由来）、
    右にナラティブ（この場面で・この画面が・こう使われる）。画像が無い場合は
    テキスト概略＋注記のみで構成する（create-design-doc 未実行時）。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, slide_spec["title"], brand)

    image = slide_spec.get("image")
    top = 1.6
    bottom_margin = 0.7
    box_h = SLIDE_H_IN - top - bottom_margin
    image_w = 5.0
    gap = 0.4
    narrative_left = 0.6 + image_w + gap
    narrative_w = SLIDE_W_IN - narrative_left - 0.6

    if image:
        add_rect(slide, 0.6, top, image_w, box_h, brand["surface"])
        add_picture_fit(slide, image, 0.6 + 0.1, top + 0.1, image_w - 0.2, box_h - 0.2)
    else:
        add_rect(slide, 0.6, top, image_w, box_h, brand["surface"])
        _, note_tf = add_textbox(slide, 0.8, top, image_w - 0.4, box_h, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraph_text(
            note_tf, slide_spec.get("note", ""), brand, size=13, color=brand["muted"],
            align=PP_ALIGN.CENTER, first=True,
        )
        narrative_left = 0.6 + image_w + gap

    _, narrative_tf = add_textbox(slide, narrative_left, top, narrative_w, box_h)
    add_paragraph_text(
        narrative_tf, slide_spec["narrative"], brand, size=16, color=brand["text"], first=True,
    )
    return slide


def build_approval(prs, slide_spec, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, brand["secondary"])
    _, title_tf = add_textbox(slide, 0.6, 0.5, SLIDE_W_IN - 1.2, 0.9)
    add_paragraph_text(title_tf, slide_spec["title"], brand, size=28, bold=True, color="#FFFFFF", first=True)
    _, tf = add_textbox(slide, 0.6, 1.6, SLIDE_W_IN - 1.2, SLIDE_H_IN - 2.3)
    items = slide_spec["items"]
    for i, item in enumerate(items):
        text = f"{i + 1}. [{item['kind']}] {item['text']}"
        add_paragraph_text(tf, text, brand, size=18, color="#FFFFFF", first=(i == 0))
        if i < len(items) - 1:
            tf.paragraphs[-1].space_after = Pt(12)
    return slide


BUILDERS = {
    "cover": None,  # 個別処理
    "premise": build_premise,
    "problem": build_bullets_slide,
    "solution": build_bullets_slide,
    "scope": build_scope,
    "winning": build_winning,
    "screen_flow": build_screen_flow,
    "validation": build_bullets_slide,
    "approval": None,  # 個別処理
}


def build_deck(spec: dict, output: Path) -> None:
    brand = {**DEFAULT_BRAND, **spec.get("brand", {})}
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slides = spec["slides"]
    total = len(slides)
    deck_title = spec.get("title", "")

    for idx, slide_spec in enumerate(slides, start=1):
        role = slide_spec["role"]
        if role == "cover":
            build_cover(prs, slide_spec, brand, spec)
            continue
        if role == "approval":
            build_approval(prs, slide_spec, brand)
            continue
        builder = BUILDERS[role]
        slide = builder(prs, slide_spec, brand)
        add_footer(slide, idx, total, brand, deck_title)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> int:
    parser = argparse.ArgumentParser(description="deck spec (JSON) から PPTX を生成する")
    parser.add_argument("--spec", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()

    try:
        spec = load_spec(args.spec)
        validate_spec(spec)
        build_deck(spec, args.output)
    except SpecError as exc:
        print(f"SpecError: {exc}", file=sys.stderr)
        return 1

    print(f"生成完了: {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
